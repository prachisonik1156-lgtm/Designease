# ─── IMPORTS ──────────────────────────────────────────────────────────────────
import sqlite3
import hashlib
import secrets
import os
import json
import re
import io
import random
from datetime import datetime, timedelta
from functools import wraps
from pathlib import Path
from io import BytesIO

from flask import (
    Flask, render_template, redirect, url_for,
    request, session, jsonify, send_from_directory, send_file, flash
)

# AI Presentation imports — optional
try:
    import anthropic
    import requests
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.pdfgen import canvas as pdf_canvas
    from reportlab.lib.utils import ImageReader
    from PIL import Image
    from pptx import Presentation as PPTPresentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    AI_ENABLED = True
except ImportError:
    AI_ENABLED = False
    print("⚠️  AI features disabled. Install: pip install anthropic requests reportlab Pillow")

# ─── CONFIG ───────────────────────────────────────────────────────────────────
ADMIN_EMAIL    = "admin@designease.com"
ADMIN_PASSWORD = "admin@1156"

BASE_DIR            = Path(__file__).parent
UPLOAD_DIR          = BASE_DIR / "static" / "uploads"
TEMPLATE_UPLOAD_DIR = BASE_DIR / "static" / "uploads" / "templates"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
TEMPLATE_UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

ALLOWED_EXT = {'png', 'jpg', 'jpeg', 'gif', 'webp', 'svg'}

app = Flask(
    __name__,
    template_folder=str(BASE_DIR / 'templates'),
    static_folder=str(BASE_DIR / 'static'),
)
app.config['UPLOAD_FOLDER']      = str(UPLOAD_DIR)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024
app.config['JSON_SORT_KEYS']     = False
app.secret_key                   = "my_project_designease_5056_secure"

# AI clients setup
if AI_ENABLED:
    ai_client           = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY", ""))
    UNSPLASH_ACCESS_KEY = os.environ.get("UNSPLASH_ACCESS_KEY", "")
    PEXELS_API_KEY      = os.environ.get("PEXELS_API_KEY", "")
else:
    UNSPLASH_ACCESS_KEY = ""
    PEXELS_API_KEY      = ""

presentations = {}
_used_urls    = {}

# ─── DATABASE ─────────────────────────────────────────────────────────────────
DATABASE = str(BASE_DIR / "designease.db")

def get_db_connection():
    conn = sqlite3.connect(DATABASE, timeout=10)
    conn.row_factory = sqlite3.Row
    return conn

get_db = get_db_connection  # alias for admin routes

def hash_pw(pw):
    return hashlib.sha256(pw.encode()).hexdigest()

hash_password = hash_pw  # alias

def init_db():
    conn = get_db_connection()
    c    = conn.cursor()

    # ── Core tables ───────────────────────────────────────────────────────────
    c.execute("""CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT,
        name TEXT,
        email TEXT UNIQUE NOT NULL,
        password TEXT NOT NULL,
        status TEXT DEFAULT 'active',
        designs_count INTEGER DEFAULT 0,
        ai_presentations INTEGER DEFAULT 0,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        last_active TIMESTAMP DEFAULT CURRENT_TIMESTAMP)""")

    c.execute("""CREATE TABLE IF NOT EXISTS templates (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        title TEXT,
        name TEXT,
        category TEXT,
        description TEXT,
        preview_image TEXT,
        status TEXT DEFAULT 'active',
        usage_count INTEGER DEFAULT 0,
        download_count INTEGER DEFAULT 0,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)""")

    c.execute("""CREATE TABLE IF NOT EXISTS designs (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        name TEXT NOT NULL DEFAULT 'Untitled Design',
        canvas_data TEXT NOT NULL DEFAULT '{}',
        thumbnail TEXT DEFAULT '',
        width INTEGER DEFAULT 1080,
        height INTEGER DEFAULT 1080,
        created_at TEXT DEFAULT (datetime('now')),
        updated_at TEXT DEFAULT (datetime('now')),
        FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE)""")

    c.execute("""CREATE TABLE IF NOT EXISTS uploads (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        filename TEXT NOT NULL,
        original_name TEXT NOT NULL,
        url TEXT NOT NULL,
        size INTEGER DEFAULT 0,
        created_at TEXT DEFAULT (datetime('now')),
        FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE)""")

    # ── Admin tables ──────────────────────────────────────────────────────────
    c.execute("""CREATE TABLE IF NOT EXISTS admins (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE NOT NULL,
        password TEXT NOT NULL,
        email TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)""")

    c.execute("""CREATE TABLE IF NOT EXISTS activity_log (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        action TEXT NOT NULL,
        details TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)""")

    c.execute("""CREATE TABLE IF NOT EXISTS system_config (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        config_key TEXT UNIQUE NOT NULL,
        config_value TEXT NOT NULL,
        description TEXT,
        key TEXT PRIMARY KEY,
        value TEXT NOT NULL,
        updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)""")
    
    c.execute("""CREATE TABLE IF NOT EXISTS system_settings (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        setting_key TEXT UNIQUE NOT NULL,
        setting_value TEXT NOT NULL)""")

    c.execute("INSERT OR IGNORE INTO system_settings (setting_key, setting_value) VALUES ('user_registration', '1')")
    c.execute("INSERT OR IGNORE INTO system_settings (setting_key, setting_value) VALUES ('maintenance_mode', '0')")
    c.execute("INSERT OR IGNORE INTO system_settings (setting_key, setting_value) VALUES ('ai_enabled', '1')")

    c.execute("""CREATE TABLE IF NOT EXISTS daily_stats (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        stat_date TEXT UNIQUE NOT NULL,
        designs_created INTEGER DEFAULT 0,
        ai_presentations INTEGER DEFAULT 0,
        new_users INTEGER DEFAULT 0,
        active_users INTEGER DEFAULT 0)""")
    
    c.execute("""CREATE TABLE IF NOT EXISTS contact_messages (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    first_name TEXT NOT NULL,
    last_name TEXT NOT NULL,
    email TEXT NOT NULL,
    phone TEXT,
    subject TEXT,
    message TEXT NOT NULL,
    status TEXT DEFAULT 'unread',
    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)""")

    # ── Migrations — safely add missing columns to existing tables ────────────
    existing = {row[1] for row in c.execute("PRAGMA table_info(templates)").fetchall()}
    for col, defn in [
        ("name",           "TEXT"),
        ("title",          "TEXT"),
        ("description",    "TEXT"),
        ("preview_image",  "TEXT"),
        ("status",         "TEXT DEFAULT 'active'"),
        ("usage_count",    "INTEGER DEFAULT 0"),
        ("download_count", "INTEGER DEFAULT 0"),
    ]:
        if col not in existing:
            c.execute(f"ALTER TABLE templates ADD COLUMN {col} {defn}")

    existing = {row[1] for row in c.execute("PRAGMA table_info(users)").fetchall()}
    for col, defn in [
        ("username",         "TEXT"),
        ("name",             "TEXT"),
        ("status",           "TEXT DEFAULT 'active'"),
        ("designs_count",    "INTEGER DEFAULT 0"),
        ("ai_presentations", "INTEGER DEFAULT 0"),
        ("last_active",      "TIMESTAMP"),
        ("created_at",       "TIMESTAMP"),
    ]:
        if col not in existing:
            c.execute(f"ALTER TABLE users ADD COLUMN {col} {defn}")

    # ── Seed admin ────────────────────────────────────────────────────────────
    if not c.execute("SELECT 1 FROM admins WHERE username='admin'").fetchone():
        c.execute("INSERT INTO admins (username, password, email) VALUES (?, ?, ?)",
                  ('admin', hash_password('admin123'), ADMIN_EMAIL))

    # ── Seed system config ────────────────────────────────────────────────────
    for key, val, desc in [
        ('ai_enabled',         'true',            'Enable/Disable AI Presentation Generator'),
        ('max_ai_per_user',    '10',              'Max AI presentations per user per month'),
        ('maintenance_mode',   'false',           'Enable maintenance mode'),
        ('allow_registration', 'true',            'Allow new user registrations'),
        ('template_approval',  'false',           'Require admin approval for templates'),
        ('ai_model',           'claude-opus-4-6', 'AI Model for presentation generation'),
    ]:
        c.execute("INSERT OR IGNORE INTO system_config (config_key, config_value, description) VALUES (?,?,?)",
                  (key, val, desc))

    # ── Seed sample users ─────────────────────────────────────────────────────
    if c.execute("SELECT COUNT(*) FROM users").fetchone()[0] == 0:
        for uname, email, designs, ai in [
            ('Alice Johnson', 'alice@email.com', 45, 12),
            ('Bob Martinez',  'bob@email.com',   32,  8),
            ('Carol White',   'carol@email.com', 67, 23),
            ('David Kim',     'david@email.com', 15,  3),
            ('Emma Davis',    'emma@email.com',  89, 31),
        ]:
            days_ago    = random.randint(1, 90)
            created     = (datetime.now() - timedelta(days=days_ago)).strftime('%Y-%m-%d %H:%M:%S')
            last_active = (datetime.now() - timedelta(days=random.randint(0, days_ago))).strftime('%Y-%m-%d %H:%M:%S')
            status      = 'blocked' if random.random() < 0.1 else 'active'
            c.execute("""INSERT OR IGNORE INTO users
                         (name, username, email, password, designs_count, ai_presentations, created_at, last_active, status)
                         VALUES (?,?,?,?,?,?,?,?,?)""",
                      (uname, uname.split()[0].lower(), email, hash_password('password'),
                       designs, ai, created, last_active, status))

    # ── Seed sample templates ─────────────────────────────────────────────────
    if c.execute("SELECT COUNT(*) FROM templates").fetchone()[0] == 0:
        for name, cat, desc, usage in [
            ('Modern Business Card', 'Poster',       'Sleek and professional business card template', 47),
            ('Creative Resume',      'Resume',       'Stand-out resume for creative professionals',   89),
            ('Startup Pitch Deck',   'Presentation', 'Investor-ready startup presentation',          134),
            ('Minimalist Poster',    'Poster',       'Clean and minimal event poster',                62),
            ('Corporate Resume',     'Resume',       'Professional corporate resume template',        73),
        ]:
            created = (datetime.now() - timedelta(days=random.randint(1, 60))).strftime('%Y-%m-%d %H:%M:%S')
            c.execute("INSERT OR IGNORE INTO templates (name, title, category, description, usage_count, created_at) VALUES (?,?,?,?,?,?)",
                      (name, name, cat, desc, usage, created))

    # ── Seed daily stats ──────────────────────────────────────────────────────
    if c.execute("SELECT COUNT(*) FROM daily_stats").fetchone()[0] == 0:
        for i in range(30):
            date = (datetime.now() - timedelta(days=i)).strftime('%Y-%m-%d')
            c.execute("INSERT OR IGNORE INTO daily_stats (stat_date,designs_created,ai_presentations,new_users,active_users) VALUES (?,?,?,?,?)",
                      (date, random.randint(20,120), random.randint(5,40), random.randint(0,8), random.randint(30,150)))

    # ── Seed activity log ─────────────────────────────────────────────────────
    if c.execute("SELECT COUNT(*) FROM activity_log").fetchone()[0] == 0:
        for action, atype in [
            ('User alice@email.com created a new design',    'design'),
            ('User bob@email.com generated AI presentation', 'ai'),
            ('Template "Modern Business Card" used 5 times', 'template'),
            ('New user emma@email.com registered',           'user'),
            ('Admin updated system configuration',           'config'),
        ]:
            created = (datetime.now() - timedelta(minutes=random.randint(5,1440))).strftime('%Y-%m-%d %H:%M:%S')
            c.execute("INSERT INTO activity_log (action, details, created_at) VALUES (?,?,?)",
                      (action, atype, created))

    conn.commit()
    conn.close()

# ─── HELPERS ──────────────────────────────────────────────────────────────────
def safe_filename(original):
    ext = original.rsplit('.', 1)[-1].lower() if '.' in original else 'bin'
    return f"{secrets.token_hex(12)}.{ext}"

def get_user_id():
    user = session.get('user')
    if not user: return None
    if isinstance(user, int): return user
    conn = get_db_connection()
    row  = conn.execute("SELECT id FROM users WHERE username=? OR email=?", (user, user)).fetchone()
    conn.close()
    return row['id'] if row else None

def login_required(f):
    @wraps(f)
    def wrapper(*args, **kwargs):
        if not session.get('user'):
            if request.is_json: return jsonify({'error': 'Login required'}), 401
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return wrapper

def admin_required(f):
    @wraps(f)
    def wrapper(*args, **kwargs):
        if not session.get('admin'):
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return wrapper

def get_setting(key, default='1'):
    conn = get_db_connection()
    row = conn.execute(
        "SELECT setting_value FROM system_settings WHERE key=?", (key,)
    ).fetchone()
    conn.close()
    return row['setting_value'] if row else default

# ─── IMAGE DATABASE ───────────────────────────────────────────────────────────
IMAGE_DB = {
    "artificial intelligence": [
        "https://images.unsplash.com/photo-1677442135703-1787eea5ce01?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1620712943543-bcc4688e7485?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1655720828018-edd2daec9349?w=800&h=500&fit=crop&auto=format",
    ],
    "machine learning": [
        "https://images.unsplash.com/photo-1555949963-aa79dcee981c?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1504639725590-34d0984388bd?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1526628652108-aa545604566e?w=800&h=500&fit=crop&auto=format",
    ],
    "neural network": [
        "https://images.unsplash.com/photo-1635070041078-e363dbe005cb?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1558494949-ef010cbdcc31?w=800&h=500&fit=crop&auto=format",
    ],
    "robot": [
        "https://images.unsplash.com/photo-1485827404703-89b55fcc595e?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1507146153580-69a1fe6d8aa1?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1531746790731-6c087fecd65a?w=800&h=500&fit=crop&auto=format",
    ],
    "computer": [
        "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1496181133206-80ce9b88a853?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1593642632559-0c6d3fc62b89?w=800&h=500&fit=crop&auto=format",
    ],
    "programming": [
        "https://images.unsplash.com/photo-1542831371-29b0f74f9713?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1461749280684-ddefd8050e31?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1518432031352-d6fc5c10da5a?w=800&h=500&fit=crop&auto=format",
    ],
    "software": [
        "https://images.unsplash.com/photo-1555066931-4365d14bab8c?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1504384308090-c894fdcc538d?w=800&h=500&fit=crop&auto=format",
    ],
    "technology": [
        "https://images.unsplash.com/photo-1518770660439-4636190af475?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1488590528505-98d2b5aba04b?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1550751827-4bd374c3f58b?w=800&h=500&fit=crop&auto=format",
    ],
    "cloud computing": [
        "https://images.unsplash.com/photo-1544197150-b99a580bb7a8?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1451187580459-43490279c0fa?w=800&h=500&fit=crop&auto=format",
    ],
    "cybersecurity": [
        "https://images.unsplash.com/photo-1550751827-4bd374c3f58b?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1563013544-824ae1b704d3?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1614064641938-3bbee52942c7?w=800&h=500&fit=crop&auto=format",
    ],
    "internet": [
        "https://images.unsplash.com/photo-1558494949-ef010cbdcc31?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1484417894907-623942c8ee29?w=800&h=500&fit=crop&auto=format",
    ],
    "smartphone": [
        "https://images.unsplash.com/photo-1511707171634-5f897ff02aa9?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1512941937669-90a1b58e7e9c?w=800&h=500&fit=crop&auto=format",
    ],
    "data": [
        "https://images.unsplash.com/photo-1551288049-bebda4e38f71?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1504868584819-f8e8b4b6d7e3?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1460925895917-adf4e566c039?w=800&h=500&fit=crop&auto=format",
    ],
    "analytics": [
        "https://images.unsplash.com/photo-1543286386-713bdd548da4?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1611974789855-9c2a0a7236a3?w=800&h=500&fit=crop&auto=format",
    ],
    "automation": [
        "https://images.unsplash.com/photo-1565043666747-69f6646db940?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1518770660439-4636190af475?w=800&h=500&fit=crop&auto=format",
    ],
    "business": [
        "https://images.unsplash.com/photo-1552664730-d307ca884978?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1507003211169-0a1dd7228f2d?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1454165804606-c3d57bc86b40?w=800&h=500&fit=crop&auto=format",
    ],
    "startup": [
        "https://images.unsplash.com/photo-1559136555-9303baea8ebd?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1497366216548-37526070297c?w=800&h=500&fit=crop&auto=format",
    ],
    "finance": [
        "https://images.unsplash.com/photo-1611974789855-9c2a0a7236a3?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1579621970563-ebec7560ff3e?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1565514020179-026b92b84bb6?w=800&h=500&fit=crop&auto=format",
    ],
    "investment": [
        "https://images.unsplash.com/photo-1559526324-593bc073d938?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1611974789855-9c2a0a7236a3?w=800&h=500&fit=crop&auto=format",
    ],
    "stock market": [
        "https://images.unsplash.com/photo-1569025743873-ea3a9ade89f9?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1611974789855-9c2a0a7236a3?w=800&h=500&fit=crop&auto=format",
    ],
    "marketing": [
        "https://images.unsplash.com/photo-1533750349088-cd871a92f312?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1460925895917-adf4e566c039?w=800&h=500&fit=crop&auto=format",
    ],
    "digital marketing": [
        "https://images.unsplash.com/photo-1432888498266-38ffec3eaf0a?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1562577309-4932fdd64cd1?w=800&h=500&fit=crop&auto=format",
    ],
    "social media": [
        "https://images.unsplash.com/photo-1611162617213-7d7a39e9b1d7?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1563986768494-4dee2763ff3f?w=800&h=500&fit=crop&auto=format",
    ],
    "ecommerce": [
        "https://images.unsplash.com/photo-1556742049-0cfed4f6a45d?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1523474253046-8cd2748b5fd2?w=800&h=500&fit=crop&auto=format",
    ],
    "leadership": [
        "https://images.unsplash.com/photo-1519389950473-47ba0277781c?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1573497620053-ea5300f94f21?w=800&h=500&fit=crop&auto=format",
    ],
    "team": [
        "https://images.unsplash.com/photo-1522071820081-009f0129c71c?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1529156069898-49953e39b3ac?w=800&h=500&fit=crop&auto=format",
    ],
    "meeting": [
        "https://images.unsplash.com/photo-1517245386807-bb43f82c33c4?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1542744173-8e7e53415bb0?w=800&h=500&fit=crop&auto=format",
    ],
    "remote work": [
        "https://images.unsplash.com/photo-1593642632559-0c6d3fc62b89?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1590402494587-44b71d7772f6?w=800&h=500&fit=crop&auto=format",
    ],
    "productivity": [
        "https://images.unsplash.com/photo-1484480974693-6ca0a78fb36b?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1507925921958-8a62f3d1a50d?w=800&h=500&fit=crop&auto=format",
    ],
    "innovation": [
        "https://images.unsplash.com/photo-1485827404703-89b55fcc595e?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1495592822108-9e6261896da8?w=800&h=500&fit=crop&auto=format",
    ],
    "strategy": [
        "https://images.unsplash.com/photo-1454165804606-c3d57bc86b40?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1581090464777-f3220bbe1b8b?w=800&h=500&fit=crop&auto=format",
    ],
    "climate change": [
        "https://images.unsplash.com/photo-1504711434969-e33886168f5c?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1569163139394-de4e4f43e4e2?w=800&h=500&fit=crop&auto=format",
    ],
    "environment": [
        "https://images.unsplash.com/photo-1542601906897-43d0e21e6c0d?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1441974231531-c6227db76b6e?w=800&h=500&fit=crop&auto=format",
    ],
    "solar energy": [
        "https://images.unsplash.com/photo-1509391366360-2e959784a276?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1497440001374-f26997328c1b?w=800&h=500&fit=crop&auto=format",
    ],
    "wind energy": [
        "https://images.unsplash.com/photo-1466611653033-459dfb3ee7ce?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1532601224476-15c79f2f7a51?w=800&h=500&fit=crop&auto=format",
    ],
    "sustainability": [
        "https://images.unsplash.com/photo-1542601906897-43d0e21e6c0d?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1518531933037-91b2f5f229cc?w=800&h=500&fit=crop&auto=format",
    ],
    "forest": [
        "https://images.unsplash.com/photo-1448375240586-882707db888b?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1542273917363-3b1817f69a2d?w=800&h=500&fit=crop&auto=format",
    ],
    "ocean": [
        "https://images.unsplash.com/photo-1505118380757-91f5f5632de0?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1518020382113-a7e8fc38eac9?w=800&h=500&fit=crop&auto=format",
    ],
    "health": [
        "https://images.unsplash.com/photo-1576091160550-2173dba999ef?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1505751172876-fa1923c5c528?w=800&h=500&fit=crop&auto=format",
    ],
    "medical": [
        "https://images.unsplash.com/photo-1559757148-5c350d0d3c56?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1579684385127-1ef15d508118?w=800&h=500&fit=crop&auto=format",
    ],
    "hospital": [
        "https://images.unsplash.com/photo-1538108149393-fbbd81895907?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1516549655169-df83a0774514?w=800&h=500&fit=crop&auto=format",
    ],
    "mental health": [
        "https://images.unsplash.com/photo-1559757175-0eb30cd8c063?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1506126613408-eca07ce68773?w=800&h=500&fit=crop&auto=format",
    ],
    "fitness": [
        "https://images.unsplash.com/photo-1571019614242-c5c5dee9f50b?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1534438327276-14e5300c3a48?w=800&h=500&fit=crop&auto=format",
    ],
    "nutrition": [
        "https://images.unsplash.com/photo-1490645935967-10de6ba17061?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1512621776951-a57141f2eefd?w=800&h=500&fit=crop&auto=format",
    ],
    "science": [
        "https://images.unsplash.com/photo-1507413245164-6160d8298b31?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1532094349884-543559bd0e52?w=800&h=500&fit=crop&auto=format",
    ],
    "space": [
        "https://images.unsplash.com/photo-1446776811953-b23d57bd21aa?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1462331940025-496dfbfc7564?w=800&h=500&fit=crop&auto=format",
    ],
    "planet": [
        "https://images.unsplash.com/photo-1614728423169-3f65fd722b7e?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1506443432602-ac2fcd6f54e0?w=800&h=500&fit=crop&auto=format",
    ],
    "chemistry": [
        "https://images.unsplash.com/photo-1532187863486-abf9dbad1b69?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1628863353691-0071c8c1874c?w=800&h=500&fit=crop&auto=format",
    ],
    "education": [
        "https://images.unsplash.com/photo-1503676260728-1c00da094a0b?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1546410531-bb4caa6b424d?w=800&h=500&fit=crop&auto=format",
    ],
    "learning": [
        "https://images.unsplash.com/photo-1456513080510-7bf3a84b82f8?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1434030216411-0b793f4b4173?w=800&h=500&fit=crop&auto=format",
    ],
    "online learning": [
        "https://images.unsplash.com/photo-1610484826967-09c5720778c7?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1571260899304-425eee4c7efc?w=800&h=500&fit=crop&auto=format",
    ],
    "blockchain": [
        "https://images.unsplash.com/photo-1639762681485-074b7f938ba0?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1644361566696-3d442b5b482a?w=800&h=500&fit=crop&auto=format",
    ],
    "cryptocurrency": [
        "https://images.unsplash.com/photo-1621761191319-c6fb62004040?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1518546305927-5a555bb7020d?w=800&h=500&fit=crop&auto=format",
    ],
    "design": [
        "https://images.unsplash.com/photo-1561070791-2526d30994b5?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1558655146-364adaf1fcc9?w=800&h=500&fit=crop&auto=format",
    ],
    "architecture": [
        "https://images.unsplash.com/photo-1486325212027-8081e485255e?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1503387762-592deb58ef4e?w=800&h=500&fit=crop&auto=format",
    ],
    "music": [
        "https://images.unsplash.com/photo-1511379938547-c1f69419868d?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1507838153414-b4b713384a76?w=800&h=500&fit=crop&auto=format",
    ],
    "travel": [
        "https://images.unsplash.com/photo-1469854523086-cc02fe5d8800?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1488085061387-422e29b40080?w=800&h=500&fit=crop&auto=format",
    ],
    "city": [
        "https://images.unsplash.com/photo-1477959858617-67f85cf4f1df?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1480714378408-67cf0d13bc1b?w=800&h=500&fit=crop&auto=format",
    ],
    "food": [
        "https://images.unsplash.com/photo-1546069901-ba9599a7e63c?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1504674900247-0877df9cc836?w=800&h=500&fit=crop&auto=format",
    ],
    "sport": [
        "https://images.unsplash.com/photo-1461896836934-ffe607ba8211?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1541534741688-6078c6bfb5c5?w=800&h=500&fit=crop&auto=format",
    ],
    "psychology": [
        "https://images.unsplash.com/photo-1559757175-0eb30cd8c063?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1606326608606-aa0b62935f2b?w=800&h=500&fit=crop&auto=format",
    ],
    "introduction": [
        "https://images.unsplash.com/photo-1542626991-cbc4e32524cc?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1557804506-669a67965ba0?w=800&h=500&fit=crop&auto=format",
    ],
    "overview": [
        "https://images.unsplash.com/photo-1454165804606-c3d57bc86b40?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1507925921958-8a62f3d1a50d?w=800&h=500&fit=crop&auto=format",
    ],
    "concept": [
        "https://images.unsplash.com/photo-1495592822108-9e6261896da8?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1580582932707-520aed937b7b?w=800&h=500&fit=crop&auto=format",
    ],
    "application": [
        "https://images.unsplash.com/photo-1551650975-87deedd944c3?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1504384308090-c894fdcc538d?w=800&h=500&fit=crop&auto=format",
    ],
    "benefit": [
        "https://images.unsplash.com/photo-1518186285589-2f7649de83e0?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1519389950473-47ba0277781c?w=800&h=500&fit=crop&auto=format",
    ],
    "conclusion": [
        "https://images.unsplash.com/photo-1499750310107-5fef28a66643?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1434030216411-0b793f4b4173?w=800&h=500&fit=crop&auto=format",
    ],
    "future": [
        "https://images.unsplash.com/photo-1485827404703-89b55fcc595e?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1451187580459-43490279c0fa?w=800&h=500&fit=crop&auto=format",
    ],
    "growth": [
        "https://images.unsplash.com/photo-1579621970795-87facc2f976d?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1460925895917-adf4e566c039?w=800&h=500&fit=crop&auto=format",
    ],
    "presentation": [
        "https://images.unsplash.com/photo-1557804506-669a67965ba0?w=800&h=500&fit=crop&auto=format",
        "https://images.unsplash.com/photo-1517245386807-bb43f82c33c4?w=800&h=500&fit=crop&auto=format",
    ],
}

# ─── IMAGE FETCH ──────────────────────────────────────────────────────────────
def fetch_image_for_slide(query, slide_index=0, presentation_id=None):
    used = _used_urls.get(presentation_id, set()) if presentation_id else set()

    if AI_ENABLED and PEXELS_API_KEY:
        try:
            r = requests.get("https://api.pexels.com/v1/search",
                headers={"Authorization": PEXELS_API_KEY},
                params={"query": query, "per_page": 10, "orientation": "landscape"}, timeout=6)
            if r.status_code == 200:
                for p in r.json().get("photos", []):
                    url = p["src"]["large"]
                    if url not in used:
                        if presentation_id: _used_urls.setdefault(presentation_id, set()).add(url)
                        return url
        except Exception as e: print(f"Pexels error: {e}")

    if AI_ENABLED and UNSPLASH_ACCESS_KEY:
        try:
            r = requests.get("https://api.unsplash.com/search/photos",
                headers={"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"},
                params={"query": query, "per_page": 10, "orientation": "landscape"}, timeout=6)
            if r.status_code == 200:
                for item in r.json().get("results", []):
                    url = item["urls"]["regular"]
                    if url not in used:
                        if presentation_id: _used_urls.setdefault(presentation_id, set()).add(url)
                        return url
        except Exception as e: print(f"Unsplash error: {e}")

    query_lower = query.lower().strip()
    words = [w for w in query_lower.split() if len(w) > 3]
    scores = {}
    for keyword, urls in IMAGE_DB.items():
        score = 0
        if keyword == query_lower: score += 20
        elif keyword in query_lower: score += 10
        elif query_lower in keyword: score += 8
        for w in words:
            if w in keyword: score += 3
            elif keyword in w: score += 2
        if score > 0: scores[keyword] = score

    for keyword, _ in sorted(scores.items(), key=lambda x: -x[1]):
        for url in IMAGE_DB[keyword]:
            if url not in used:
                if presentation_id: _used_urls.setdefault(presentation_id, set()).add(url)
                return url

    all_urls = []
    for urls in IMAGE_DB.values(): all_urls.extend(urls)
    seen = set()
    unique_urls = [u for u in all_urls if not (u in seen or seen.add(u))]
    rotated = unique_urls[slide_index:] + unique_urls[:slide_index]
    for url in rotated:
        if url not in used:
            if presentation_id: _used_urls.setdefault(presentation_id, set()).add(url)
            return url
    return unique_urls[slide_index % len(unique_urls)]

# ─── SLIDE GENERATION ─────────────────────────────────────────────────────────
def generate_presentation_slides(topic, presentation_id):
    if not AI_ENABLED: return get_fallback_slides(topic, presentation_id)
    try:
        prompt = f"""Create a professional 7-slide presentation about "{topic}".

Return ONLY valid JSON (no markdown, no code blocks):
{{
  "slides": [
    {{
      "title": "Slide title here",
      "content": ["Complete unique sentence 1 (15-20 words).", "Sentence 2.", "Sentence 3.", "Sentence 4.", "Sentence 5."],
      "imageKeywords": "specific visual keyword for THIS slide only"
    }}
  ]
}}

SLIDE STRUCTURE for "{topic}":
- Slide 1: What it is and why it matters
- Slide 2: Historical background and origin
- Slide 3: Key concepts and terminology
- Slide 4: Features and characteristics
- Slide 5: Real-world applications
- Slide 6: Benefits and advantages
- Slide 7: Conclusion and future

imageKeywords: Every slide MUST have a DIFFERENT keyword (2-4 words, concrete visual terms).
Each slide content must be COMPLETELY DIFFERENT. No repeating sentences."""

        message = ai_client.messages.create(
            model="claude-opus-4-6", max_tokens=4000,
            messages=[{"role": "user", "content": prompt}])
        response_text = message.content[0].text
        if "```json" in response_text:
            response_text = response_text.split("```json")[1].split("```")[0].strip()
        elif "```" in response_text:
            response_text = response_text.split("```")[1].split("```")[0].strip()
        slides = json.loads(response_text)["slides"]
        _used_urls[presentation_id] = set()
        for i, slide in enumerate(slides):
            kw = slide.get("imageKeywords", f"{topic} slide {i+1}")
            slide["imageUrl"] = fetch_image_for_slide(kw, i, presentation_id)
        return slides
    except Exception as e:
        print(f"Slide generation error: {e}")
        return get_fallback_slides(topic, presentation_id)

def get_fallback_slides(topic, presentation_id=None):
    _used_urls[presentation_id] = set()
    configs = [
        (f"{topic}: Overview",         topic),
        ("Introduction",                f"introduction context {topic}"),
        ("Key Concepts",                f"{topic} technology concept"),
        ("Features & Characteristics",  f"{topic} features detail"),
        ("Real-World Applications",     f"{topic} application industry"),
        ("Benefits & Impact",           f"benefit growth success {topic}"),
        ("Conclusion & Future",         f"future outlook {topic}"),
    ]
    all_content = [
        [f"{topic} is a rapidly evolving field with significant real-world impact across industries.",
         f"This presentation covers core concepts, features, applications, and future outlook of {topic}.",
         f"Understanding {topic} is essential for professionals, researchers, and decision-makers today.",
         "We explore both theoretical foundations and practical implementations throughout these slides.",
         f"Each section provides structured, detailed insights to help you grasp {topic} comprehensively."],
        [f"{topic} has emerged as one of the most transformative areas in modern technology and society.",
         "Its development accelerated due to increased data availability and computing power.",
         "Key drivers include global connectivity, innovation demand, and cross-industry collaboration.",
         f"Today, {topic} influences how businesses operate, make decisions, and stay competitive.",
         "Understanding its origins helps appreciate its current capabilities and future direction."],
        [f"The foundation of {topic} rests on well-defined principles, models, and structured frameworks.",
         "Core components are interdependent — understanding one provides insight into all others.",
         "Researchers continuously refine these concepts based on empirical evidence and real-world results.",
         "Fundamental theories provide the basis for developing advanced tools and specialized systems.",
         "Mastery of these concepts is critical before exploring higher-level applications or extensions."],
        [f"{topic} offers distinct features that set it apart from conventional or traditional methods.",
         "Scalability allows solutions to grow efficiently as usage, data volume, and complexity increase.",
         "Adaptability ensures continued relevance across changing environments, industries, and requirements.",
         "Key attributes include automation, speed, accuracy, and consistently measurable performance gains.",
         f"These combined features make {topic} a powerful competitive advantage in modern organizations."],
        [f"{topic} is actively deployed in healthcare, finance, education, manufacturing, and logistics.",
         "In healthcare, it improves diagnostics accuracy, treatment planning, and patient outcome tracking.",
         "Financial institutions leverage it for fraud detection, risk modeling, and customer personalization.",
         f"Educational platforms use {topic} to deliver adaptive content and personalized learning pathways.",
         "Manufacturers benefit through predictive maintenance, supply chain optimization, and quality control."],
        [f"Organizations adopting {topic} see measurable improvements in efficiency, output, and cost reduction.",
         "Decision-making becomes faster and more accurate with real-time data analysis and automation.",
         "It significantly reduces human error by automating repetitive, complex, and data-intensive tasks.",
         "Long-term benefits include stronger competitive positioning, innovation capacity, and scalable growth.",
         f"Teams can focus on creative work while {topic} handles operational complexity."],
        [f"{topic} is a transformative force fundamentally reshaping industries worldwide.",
         "Organizations that invest and adapt early will gain lasting advantages over slower competitors.",
         "Future advancements will bring greater capabilities, deeper integration, and broader societal impact.",
         f"Ethical considerations and governance will shape how {topic} evolves ahead.",
         "Staying informed, building expertise, and embracing change are keys to long-term success."],
    ]
    return [{"title": title, "content": all_content[i],
             "imageUrl": fetch_image_for_slide(kw, slide_index=i, presentation_id=presentation_id),
             "imageKeywords": kw}
            for i, (title, kw) in enumerate(configs)]


# =============================================================================
#  AUTH ROUTES
# =============================================================================

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/start")
def start():
    if "user" in session: return redirect(url_for("dashboard"))
    return redirect(url_for("signup"))

@app.route("/login", methods=["GET", "POST"])
def login():
    if get_setting('maintenance_mode') == '1':
        return render_template('maintenance.html'), 503
    
    if request.method == "POST":
        email    = request.form.get("email", "").strip()
        password = request.form.get("password", "").strip()
        if not email or not password:
            return "Email or Password missing"
        if email == ADMIN_EMAIL and password == ADMIN_PASSWORD:
            session.clear()
            session["admin"]           = True
            session["admin_logged_in"] = True
            session["admin_username"]  = "Admin"
            session["_login_flash"]    = "admin"
            return redirect(url_for("admin_dashboard"))
        conn   = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("SELECT * FROM users WHERE email=? AND password=?", (email, password))
        user = cursor.fetchone(); conn.close()
        if user:
            session.clear()
            # username column index 1, fallback to email if username empty
            session["user"] = user["username"] or user["name"] or user["email"]
            session["_login_flash"] = "user"
            return redirect(url_for("dashboard"))
        flash("Invalid email or password. Please try again.", "error")
        return redirect(url_for('login'))  
    return render_template("login.html") 

@app.route("/signup", methods=["GET", "POST"])
def signup():
    if get_setting('user_registration') == '0':
        flash('New registrations are currently disabled.', 'error')
        return redirect(url_for('login'))
    
    if request.method == "POST":
        conn = None
        try:
            username = request.form["username"]
            email    = request.form["email"]
            password = request.form["password"]
            if len(password) < 6:
                return render_template("signup.html", error="Password must be at least 6 characters.")
            conn = get_db_connection()
            conn.execute("INSERT INTO users (username, name, email, password) VALUES (?,?,?,?)",
                         (username, username, email, password))
            conn.commit()
            session["user"] = username
            session["_login_flash"] = "signup"
            return redirect(url_for("dashboard"))
        except sqlite3.IntegrityError:
            return render_template("signup.html", error="This email is already registered. Please login.")
        except Exception as e:
            return render_template("signup.html", error=f"Something went wrong. Please try again.")
        finally:
            if conn: conn.close()
    return render_template("signup.html")


@app.route("/forgot")
def forgot(): return render_template("forgot.html")

@app.route("/reset-password", methods=["POST"])
def reset_password():
    try:
        data     = request.json or {}
        email    = (data.get("email") or "").strip()
        password = (data.get("password") or "").strip()

        if not email or not password:
            return jsonify({"success": False, "error": "Email and password are required."})
        if len(password) < 6:
            return jsonify({"success": False, "error": "Password must be at least 6 characters."})

        conn = get_db_connection()
        user = conn.execute("SELECT id FROM users WHERE email=?", (email,)).fetchone()
        if not user:
            conn.close()
            return jsonify({"success": False, "error": "No account found with this email."})

        conn.execute("UPDATE users SET password=? WHERE email=?", (password, email))
        conn.commit()
        conn.close()
        return jsonify({"success": True})
    except Exception as e:
        print(f"Reset password error: {e}")
        return jsonify({"success": False, "error": "Server error. Please try again."})

@app.route("/logout")
def logout():
    flash("Logged out successfully!", "logout")
    session.pop("user", None)
    session.pop("admin", None)
    session.pop("admin_logged_in", None)
    session.pop("admin_username", None)
    return redirect(url_for("login"))

# =============================================================================
#  USER PAGES
# =============================================================================

@app.route("/clear-login-flash", methods=["POST"])
def clear_login_flash():
    session.pop("_login_flash", None)
    return "", 204

@app.route("/dashboard")
def dashboard():
    if not session.get("user"): return redirect(url_for("login"))
    return render_template("dashboard.html")

@app.route('/api/user-info')
def user_info():
    if not session.get('user'):
        return jsonify({'error': 'Not logged in'}), 401
    uid = get_user_id()
    conn = get_db_connection()
    row = conn.execute(
        "SELECT username, email FROM users WHERE id=?", (uid,)
    ).fetchone()
    conn.close()
    if not row:
        return jsonify({'error': 'User not found'}), 404
    return jsonify({
        'username': row['username'] or session.get('user'),
        'email': row['email']
    })

@app.route("/about")
def about(): return render_template("about.html")

@app.route("/contact")
def contact(): return render_template("contact.html")

# =============================================================================
#  ADMIN ROUTES
# =============================================================================

@app.route("/admin")
@admin_required
def admin_dashboard():
    conn = get_db()
    total_users     = conn.execute("SELECT COUNT(*) FROM users").fetchone()[0]
    total_templates = len(EDITOR_TEMPLATE_LIST)
    total_designs   = conn.execute("SELECT SUM(designs_count) FROM users").fetchone()[0] or 0
    total_ai        = conn.execute("SELECT SUM(ai_presentations) FROM users").fetchone()[0] or 0
    active_users    = conn.execute("SELECT COUNT(*) FROM users WHERE status='active'").fetchone()[0]
    recent_activity = conn.execute("SELECT * FROM activity_log ORDER BY created_at DESC LIMIT 8").fetchall()
    recent_users    = conn.execute("SELECT * FROM users ORDER BY created_at DESC LIMIT 5").fetchall()

    chart_labels, chart_designs, chart_ai = [], [], []
    for i in range(6, -1, -1):
        date  = (datetime.now() - timedelta(days=i)).strftime('%Y-%m-%d')
        label = (datetime.now() - timedelta(days=i)).strftime('%b %d')
        stat  = conn.execute("SELECT * FROM daily_stats WHERE stat_date=?", (date,)).fetchone()
        chart_labels.append(label)
        chart_designs.append(stat['designs_created'] if stat else 0)
        chart_ai.append(stat['ai_presentations'] if stat else 0)

    from collections import Counter
    cats       = Counter(t['cat'] for t in EDITOR_TEMPLATE_LIST)
    cat_labels = list(cats.keys())
    cat_counts = list(cats.values())
    conn.close()

    return render_template('admin_dashboard.html',
        total_users=total_users, total_templates=total_templates,
        total_designs=total_designs, total_ai=total_ai,
        active_users=active_users, recent_activity=recent_activity,
        recent_users=recent_users,
        chart_labels=json.dumps(chart_labels),
        chart_designs=json.dumps(chart_designs),
        chart_ai=json.dumps(chart_ai),
        cat_labels=json.dumps(cat_labels),
        cat_counts=json.dumps(cat_counts))

@app.route('/api/contact', methods=['POST'])
def submit_contact():
    try:
        data = request.json or {}
        first_name = data.get('first_name', '').strip()
        last_name  = data.get('last_name', '').strip()
        email      = data.get('email', '').strip()
        phone      = data.get('phone', '').strip()
        subject    = data.get('subject', '').strip()
        message    = data.get('message', '').strip()

        if not first_name or not email or not message:
            return jsonify({'error': 'Name, email aur message required hai'}), 400

        conn = get_db_connection()
        conn.execute("""INSERT INTO contact_messages
            (first_name, last_name, email, phone, subject, message)
            VALUES (?,?,?,?,?,?)""",
            (first_name, last_name, email, phone, subject, message))
        conn.execute("INSERT INTO activity_log (action, details) VALUES (?,?)",
            (f'New contact message from {first_name} {last_name} ({email})', 'contact'))
        conn.commit()
        conn.close()
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Admin messages route:
@app.route('/admin/messages')
@admin_required
def admin_messages():
    conn = get_db()
    messages = conn.execute(
        "SELECT * FROM contact_messages ORDER BY created_at DESC"
    ).fetchall()
    unread = conn.execute(
        "SELECT COUNT(*) FROM contact_messages WHERE status='unread'"
    ).fetchone()[0]
    conn.close()
    return render_template('admin_messages.html', messages=messages, unread=unread)

@app.route('/admin/messages/read/<int:mid>', methods=['POST'])
@admin_required
def mark_message_read(mid):
    conn = get_db()
    conn.execute("UPDATE contact_messages SET status='read' WHERE id=?", (mid,))
    conn.commit()
    conn.close()
    return jsonify({'success': True})

@app.route('/admin/messages/delete/<int:mid>', methods=['POST'])
@admin_required
def delete_message(mid):
    conn = get_db()
    conn.execute("DELETE FROM contact_messages WHERE id=?", (mid,))
    conn.commit()
    conn.close()
    return jsonify({'success': True})

# ── Manage Templates ──────────────────────────────────────────────────────────
@app.route("/admin/manage-templates")
@admin_required
def manage_templates():
    from collections import Counter
    cats = Counter(t['cat'] for t in EDITOR_TEMPLATE_LIST)
    cat_counts = {
        'poster':       cats.get('poster', 0),
        'resume':       cats.get('resume', 0),
        'presentation': cats.get('presentation', 0),
        'social':       cats.get('social', 0),
        'travel':       cats.get('travel', 0),
    }
    return render_template("manage_templates.html",
                           templates=EDITOR_TEMPLATE_LIST,
                           total_templates=len(EDITOR_TEMPLATE_LIST),
                           cat_counts=cat_counts)

@app.route('/admin/templates/add', methods=['POST'])
@admin_required
def add_template():
    name        = request.form.get('name','').strip()
    category    = request.form.get('category','')
    description = request.form.get('description','')
    preview_image = ''
    if 'preview_image' in request.files:
        file = request.files['preview_image']
        if file and file.filename:
            filename = f"{datetime.now().strftime('%Y%m%d%H%M%S')}_{safe_filename(file.filename)}"
            file.save(str(TEMPLATE_UPLOAD_DIR / filename))
            preview_image = filename
    conn = get_db()
    conn.execute("INSERT INTO templates (name, title, category, description, preview_image) VALUES (?,?,?,?,?)",
                 (name, name, category, description, preview_image))
    conn.execute("INSERT INTO activity_log (action, details) VALUES (?,?)",
                 (f'Template "{name}" added by admin', 'template'))
    conn.commit(); conn.close()
    flash('Template added successfully!', 'success')
    return redirect(url_for('manage_templates'))

@app.route('/admin/templates/update/<int:tid>', methods=['POST'])
@admin_required
def update_template(tid):
    name     = request.form.get('name','').strip()
    category = request.form.get('category','')
    desc     = request.form.get('description','')
    status   = request.form.get('status', 'active')
    conn = get_db()
    conn.execute("UPDATE templates SET name=?,title=?,category=?,description=?,status=? WHERE id=?",
                 (name, name, category, desc, status, tid))
    conn.execute("INSERT INTO activity_log (action, details) VALUES (?,?)",
                 (f'Template "{name}" updated by admin', 'template'))
    conn.commit(); conn.close()
    flash('Template updated!', 'success')
    return redirect(url_for('manage_templates'))

@app.route('/admin/templates/delete/<int:tid>', methods=['POST'])
@admin_required
def delete_template(tid):
    conn = get_db()
    t = conn.execute("SELECT name FROM templates WHERE id=?", (tid,)).fetchone()
    conn.execute("DELETE FROM templates WHERE id=?", (tid,))
    conn.execute("INSERT INTO activity_log (action, details) VALUES (?,?)",
                 (f'Template "{t["name"] if t else tid}" deleted', 'template'))
    conn.commit(); conn.close()
    flash('Template deleted.', 'success')
    return redirect(url_for('manage_templates'))

# ── Usage / User Reports — both old and new URL work ─────────────────────────
@app.route("/admin/user-report")
@app.route("/admin/usage-reports")
@admin_required
def user_report():
    conn = get_db()
    labels_30, designs_30, ai_30, users_30 = [], [], [], []
    for i in range(29, -1, -1):
        date  = (datetime.now() - timedelta(days=i)).strftime('%Y-%m-%d')
        label = (datetime.now() - timedelta(days=i)).strftime('%b %d')
        stat  = conn.execute("SELECT * FROM daily_stats WHERE stat_date=?", (date,)).fetchone()
        labels_30.append(label)
        designs_30.append(stat['designs_created'] if stat else 0)
        ai_30.append(stat['ai_presentations'] if stat else 0)
        users_30.append(stat['new_users'] if stat else 0)

    top_templates = conn.execute(
        "SELECT name, category, usage_count FROM templates ORDER BY usage_count DESC LIMIT 8").fetchall()

    weekly_labels, weekly_designs, weekly_ai = [], [], []
    for w in range(3, -1, -1):
        ws = datetime.now() - timedelta(weeks=w+1)
        we = datetime.now() - timedelta(weeks=w)
        d  = conn.execute(
            "SELECT SUM(designs_created), SUM(ai_presentations) FROM daily_stats WHERE stat_date BETWEEN ? AND ?",
            (ws.strftime('%Y-%m-%d'), we.strftime('%Y-%m-%d'))).fetchone()
        weekly_labels.append(f"Week {4-w}")
        weekly_designs.append(d[0] or 0)
        weekly_ai.append(d[1] or 0)

    total_users = conn.execute("SELECT COUNT(*) FROM users").fetchone()[0]
    conn.close()

    return render_template('user_report.html',
        total_users=total_users,
        labels_30=json.dumps(labels_30), designs_30=json.dumps(designs_30),
        ai_30=json.dumps(ai_30), users_30=json.dumps(users_30),
        top_templates=top_templates,
        weekly_labels=json.dumps(weekly_labels),
        weekly_designs=json.dumps(weekly_designs), weekly_ai=json.dumps(weekly_ai),
        total_designs_month=sum(designs_30), total_ai_month=sum(ai_30),
        total_new_users=sum(users_30))

# ── AI Features ───────────────────────────────────────────────────────────────
@app.route("/admin/ai-features")
@admin_required
def ai_features():
    conn = get_db()
    configs  = {row['config_key']: row['config_value']
                for row in conn.execute("SELECT * FROM system_config").fetchall()}
    total_ai = conn.execute("SELECT SUM(ai_presentations) FROM users").fetchone()[0] or 0
    ai_users = conn.execute("SELECT COUNT(*) FROM users WHERE ai_presentations > 0").fetchone()[0]
    conn.close()
    return render_template('ai_features.html', configs=configs,
                           total_ai=total_ai, ai_users=ai_users, ai_enabled=AI_ENABLED)

@app.route('/admin/ai-features/update', methods=['POST'])
@admin_required
def update_ai_features():
    ai_enabled = 'true' if request.form.get('ai_enabled') else 'false'
    max_ai     = request.form.get('max_ai_per_user', '10')
    ai_model   = request.form.get('ai_model', 'claude-opus-4-6')
    conn = get_db()
    conn.execute("UPDATE system_config SET config_value=? WHERE config_key='ai_enabled'",      (ai_enabled,))
    conn.execute("UPDATE system_config SET config_value=? WHERE config_key='max_ai_per_user'", (max_ai,))
    conn.execute("UPDATE system_config SET config_value=? WHERE config_key='ai_model'",        (ai_model,))
    conn.execute("INSERT INTO activity_log (action, details) VALUES (?,?)",
                 ('Admin updated AI feature configuration', 'config'))
    conn.commit(); conn.close()
    flash('AI settings updated!', 'success')
    return redirect(url_for('ai_features'))

# ── User Management — both old and new URL work ───────────────────────────────
@app.route("/admin/users")
@app.route("/admin/user-management")
@admin_required
def admin_users():
    conn    = get_db()
    users   = conn.execute("SELECT * FROM users ORDER BY created_at DESC").fetchall()
    total   = conn.execute("SELECT COUNT(*) FROM users").fetchone()[0]
    active  = conn.execute("SELECT COUNT(*) FROM users WHERE status='active'").fetchone()[0]
    blocked = conn.execute("SELECT COUNT(*) FROM users WHERE status='blocked'").fetchone()[0]
    conn.close()
    return render_template('admin_users.html', users=users, total=total, active=active, blocked=blocked)

@app.route('/admin/users/toggle/<int:uid>', methods=['POST'])
@admin_required
def toggle_user(uid):
    conn = get_db()
    user = conn.execute("SELECT * FROM users WHERE id=?", (uid,)).fetchone()
    new_status = 'blocked' if user['status'] == 'active' else 'active'
    conn.execute("UPDATE users SET status=? WHERE id=?", (new_status, uid))
    conn.execute("INSERT INTO activity_log (action, details) VALUES (?,?)",
                 (f'User {user["email"]} status changed to {new_status}', 'moderation'))
    conn.commit(); conn.close()
    return jsonify({"success": True, "new_status": new_status})

@app.route('/admin/users/block/<int:uid>', methods=['POST'])
@admin_required
def block_user(uid):
    conn = get_db()
    conn.execute("UPDATE users SET status='blocked' WHERE id=?", (uid,))
    conn.commit(); conn.close()
    return jsonify({"success": True})

@app.route('/admin/users/delete/<int:uid>', methods=['POST'])
@admin_required
def delete_user(uid):
    conn = get_db()
    conn.execute("DELETE FROM users WHERE id=?", (uid,))
    conn.commit(); conn.close()
    return jsonify({"success": True})

# ── System Config ─────────────────────────────────────────────────────────────
@app.route("/admin/system-config")
@admin_required
def system_config():
    conn     = get_db()
    configs  = conn.execute("SELECT * FROM system_config ORDER BY config_key").fetchall()
    activity = conn.execute("SELECT * FROM activity_log ORDER BY created_at DESC LIMIT 20").fetchall()
    total_users     = conn.execute("SELECT COUNT(*) FROM users").fetchone()[0]
    total_templates = conn.execute("SELECT COUNT(*) FROM templates").fetchone()[0]
    total_logs      = conn.execute("SELECT COUNT(*) FROM activity_log").fetchone()[0]
    conn.close()
    return render_template('system_config.html', configs=configs, activity=activity,
                           total_users=total_users, total_templates=total_templates, total_logs=total_logs)

@app.route('/admin/system-config/update', methods=['POST'])
@admin_required
def update_system_config():
    maintenance   = 'true' if request.form.get('maintenance_mode') else 'false'
    allow_reg     = 'true' if request.form.get('allow_registration') else 'false'
    tmpl_approval = 'true' if request.form.get('template_approval') else 'false'
    conn = get_db()
    conn.execute("UPDATE system_config SET config_value=? WHERE config_key='maintenance_mode'",   (maintenance,))
    conn.execute("UPDATE system_config SET config_value=? WHERE config_key='allow_registration'", (allow_reg,))
    conn.execute("UPDATE system_config SET config_value=? WHERE config_key='template_approval'",  (tmpl_approval,))
    conn.execute("INSERT INTO activity_log (action, details) VALUES (?,?)",
                 ('Admin updated system configuration settings', 'config'))
    conn.commit(); conn.close()
    flash('System configuration updated!', 'success')
    return redirect(url_for('system_config'))

# ── Admin API stats ───────────────────────────────────────────────────────────
@app.route('/admin/api/stats')
@admin_required
def admin_api_stats():
    conn  = get_db()
    today = datetime.now().strftime('%Y-%m-%d')
    stat  = conn.execute("SELECT * FROM daily_stats WHERE stat_date=?", (today,)).fetchone()
    conn.close()
    return jsonify({
        'designs_today': stat['designs_created'] if stat else 0,
        'ai_today':      stat['ai_presentations'] if stat else 0,
    })

# ── Template download ─────────────────────────────────────────────────────────
@app.route("/download/<int:template_id>")
def download_template(template_id):
    if not session.get("user"): return redirect(url_for("login"))
    conn = get_db_connection()
    conn.execute("UPDATE templates SET download_count = download_count + 1 WHERE id=?", (template_id,))
    conn.commit(); conn.close()
    return "Download started"

# =============================================================================
#  EDITOR ROUTES
# =============================================================================

@app.route('/editor')
@app.route('/editor/<int:design_id>')
def editor(design_id=None):
    if not session.get('user'):
        return redirect(url_for('login'))
    template_id = request.args.get('template', '')
    return render_template('editor.html', 
                           design_id=design_id or '', 
                           template_id=template_id)

@app.route('/create')
def create_design_redirect():
    if not session.get('user'): return redirect(url_for('login'))
    user_id = session.get('user')
    conn    = get_db_connection()
    cursor  = conn.cursor()
    if isinstance(user_id, int):
        uid = user_id
    else:
        cursor.execute("SELECT id FROM users WHERE username=? OR email=?", (user_id, user_id))
        row = cursor.fetchone()
        uid = row['id'] if row else None
    if not uid: conn.close(); return redirect(url_for('login'))
    cursor.execute("INSERT INTO designs (user_id, name, canvas_data, width, height) VALUES (?,?,?,?,?)",
                   (uid, 'Untitled Design', '{}', 1080, 1080))
    conn.commit()
    new_id = cursor.lastrowid; conn.close()
    return redirect(url_for('editor', design_id=new_id))

# =============================================================================
#  DESIGNS API
# =============================================================================

@app.route('/api/designs', methods=['GET'])
def list_designs():
    if not session.get('user'): return jsonify({'error': 'Login required'}), 401
    uid = get_user_id()
    if not uid: return jsonify({'error': 'User not found'}), 404
    conn = get_db_connection()
    rows = conn.execute(
        "SELECT id,name,thumbnail,width,height,created_at,updated_at FROM designs WHERE user_id=? ORDER BY updated_at DESC",
        (uid,)).fetchall()
    conn.close()
    return jsonify({'designs': [dict(r) for r in rows]})

@app.route('/api/designs/<int:did>', methods=['GET'])
def get_design(did):
    if not session.get('user'): return jsonify({'error': 'Login required'}), 401
    uid = get_user_id()
    conn = get_db_connection()
    row = conn.execute("SELECT * FROM designs WHERE id=? AND user_id=?", (did, uid)).fetchone()
    conn.close()
    if not row: return jsonify({'error': 'Design not found'}), 404
    return jsonify(dict(row))

@app.route('/api/designs/<int:did>', methods=['PUT'])
def save_design(did):
    if not session.get('user'): return jsonify({'error': 'Login required'}), 401
    uid = get_user_id(); d = request.json or {}
    canvas_data = json.dumps(d.get('canvas_data', {}))
    name = d.get('name', '').strip(); thumbnail = d.get('thumbnail', '')
    conn = get_db_connection()
    row = conn.execute("SELECT id FROM designs WHERE id=? AND user_id=?", (did, uid)).fetchone()
    if not row: conn.close(); return jsonify({'error': 'Permission denied'}), 403
    if name:
        conn.execute("UPDATE designs SET canvas_data=?,name=?,thumbnail=?,updated_at=datetime('now') WHERE id=?",
                     (canvas_data, name, thumbnail, did))
    else:
        conn.execute("UPDATE designs SET canvas_data=?,thumbnail=?,updated_at=datetime('now') WHERE id=?",
                     (canvas_data, thumbnail, did))
    conn.commit(); conn.close()
    return jsonify({'success': True, 'saved_at': datetime.now().isoformat()})

@app.route('/api/designs/<int:did>', methods=['DELETE'])
def delete_design(did):
    if not session.get('user'): return jsonify({'error': 'Login required'}), 401
    uid = get_user_id()
    conn = get_db_connection()
    conn.execute("DELETE FROM designs WHERE id=? AND user_id=?", (did, uid))
    conn.commit(); conn.close()
    return jsonify({'success': True})

@app.route('/api/designs/<int:did>/rename', methods=['PATCH'])
def rename_design(did):
    if not session.get('user'): return jsonify({'error': 'Login required'}), 401
    uid = get_user_id(); name = (request.json or {}).get('name', '').strip()
    if not name: return jsonify({'error': 'Name required'}), 400
    conn = get_db_connection()
    conn.execute("UPDATE designs SET name=?,updated_at=datetime('now') WHERE id=? AND user_id=?",
                 (name, did, uid))
    conn.commit(); conn.close()
    return jsonify({'success': True})

# =============================================================================
#  UPLOADS API
# =============================================================================

@app.route('/api/uploads', methods=['POST'])
def upload_image():
    if not session.get('user'): return jsonify({'error': 'Login required'}), 401
    if 'file' not in request.files: return jsonify({'error': 'No file'}), 400
    file = request.files['file']
    if not file or not file.filename: return jsonify({'error': 'Select a file'}), 400
    ext = file.filename.rsplit('.', 1)[-1].lower() if '.' in file.filename else ''
    if ext not in ALLOWED_EXT: return jsonify({'error': f'Allowed: {", ".join(ALLOWED_EXT)}'}), 400
    uid = get_user_id(); fname = safe_filename(file.filename)
    fpath = UPLOAD_DIR / fname; file.save(str(fpath))
    size = os.path.getsize(str(fpath)); url = f'/static/uploads/{fname}'
    conn = get_db_connection()
    conn.execute("INSERT INTO uploads (user_id,filename,original_name,url,size) VALUES (?,?,?,?,?)",
                 (uid, fname, file.filename, url, size))
    conn.commit(); conn.close()
    return jsonify({'success': True, 'url': url, 'filename': fname})

@app.route('/api/uploads', methods=['GET'])
def list_uploads():
    if not session.get('user'): return jsonify({'error': 'Login required'}), 401
    uid = get_user_id()
    conn = get_db_connection()
    rows = conn.execute(
        "SELECT id,filename,original_name,url,size,created_at FROM uploads WHERE user_id=? ORDER BY created_at DESC LIMIT 60",
        (uid,)).fetchall()
    conn.close()
    return jsonify({'uploads': [dict(r) for r in rows]})

@app.route('/api/uploads/<int:uid_>', methods=['DELETE'])
def delete_upload(uid_):
    if not session.get('user'): return jsonify({'error': 'Login required'}), 401
    uid = get_user_id()
    conn = get_db_connection()
    row = conn.execute("SELECT filename FROM uploads WHERE id=? AND user_id=?", (uid_, uid)).fetchone()
    if not row: conn.close(); return jsonify({'error': 'File not found'}), 404
    f = UPLOAD_DIR / row['filename']
    if f.exists(): f.unlink()
    conn.execute("DELETE FROM uploads WHERE id=?", (uid_,))
    conn.commit(); conn.close()
    return jsonify({'success': True})

@app.route('/static/uploads/<path:filename>')
def serve_upload(filename):
    return send_from_directory(str(UPLOAD_DIR), filename)

# =============================================================================
#  EDITOR TEMPLATES
# =============================================================================

EDITOR_TEMPLATE_LIST = [
    # ── Posters ──
    {'id': 'poster1',       'name': 'Hiking Poster',        'cat': 'poster',       'thumb_bg': 'linear-gradient(135deg,#ead7b5,#f5e9d6)'},
    {'id': 'poster2',       'name': 'Business Workshop',    'cat': 'poster',       'thumb_bg': 'linear-gradient(135deg,#1a1a2e,#16213e)'},
    {'id': 'poster3',       'name': 'Poster 3',             'cat': 'poster',       'thumb_bg': 'linear-gradient(135deg,#f093fb,#f5576c)'},
    {'id': 'poster4',       'name': 'Poster 4',             'cat': 'poster',       'thumb_bg': 'linear-gradient(135deg,#4facfe,#00f2fe)'},
    # ── Resumes ──
    {'id': 'resume1',       'name': 'Resume Classic',       'cat': 'resume',       'thumb_bg': 'linear-gradient(135deg,#1f2937,#374151)'},
    {'id': 'resume2',       'name': 'Resume Modern',        'cat': 'resume',       'thumb_bg': 'linear-gradient(135deg,#667eea,#764ba2)'},
    {'id': 'resume3',       'name': 'Resume Creative',      'cat': 'resume',       'thumb_bg': 'linear-gradient(135deg,#f7971e,#ffd200)'},
    {'id': 'resume4',       'name': 'Resume Minimal',       'cat': 'resume',       'thumb_bg': 'linear-gradient(135deg,#2d3748,#4a5568)'},
    {'id': 'index45',       'name': 'Resume with Photo',    'cat': 'resume',       'thumb_bg': 'linear-gradient(135deg,#e2e8f0,#f7fafc)'},
    # ── Presentations ──
    {'id': 'presentation1', 'name': 'Corporate Dark',       'cat': 'presentation', 'thumb_bg': 'linear-gradient(135deg,#1a1a2e,#0f3460)'},
    # ── Social Media ──
    {'id': 'social1',       'name': 'Food Instagram Post',  'cat': 'social',       'thumb_bg': 'linear-gradient(135deg,#f59e0b,#ef4444)'},
    {'id': 'social2',       'name': 'Quote Post Dark',      'cat': 'social',       'thumb_bg': 'linear-gradient(135deg,#0f0c29,#302b63)'},
    # ── Travel ──
    {'id': 'travel1',       'name': 'Destination Card',     'cat': 'travel',       'thumb_bg': 'linear-gradient(135deg,#0369a1,#0d9488)'},
    {'id': 'travel2',       'name': 'Tour Package',         'cat': 'travel',       'thumb_bg': 'linear-gradient(135deg,#2aa6a6,#1e8a8a)'},
]

@app.route('/api/templates')
def list_templates():
    return jsonify({'templates': EDITOR_TEMPLATE_LIST})

@app.route('/template/<template_id>')
def view_template(template_id):
    if not re.match(r'^[a-zA-Z0-9_]+$', template_id): return "Invalid template", 400
    return render_template(f'{template_id}.html')

# =============================================================================
#  AI PRESENTATION ROUTES
# =============================================================================

@app.route('/generate', methods=['GET', 'POST'])
def generate():
    if request.method == 'GET':
        if not session.get('user'): return redirect(url_for('login'))
        return render_template('generate.html')
    try:
        data  = request.json
        topic = (data.get('topic') or 'Technology').strip()
        if not topic: return jsonify({"error": "Topic is required"}), 400
        pid    = datetime.now().strftime("%Y%m%d%H%M%S%f")
        slides = generate_presentation_slides(topic, pid)
        presentations[pid] = {"id": pid, "topic": topic,
                               "created_at": datetime.now().isoformat(), "slides": slides}
        return jsonify({"success": True, "presentation_id": pid, "topic": topic, "slides": slides})
    except Exception as e:
        print(f"Generate error: {e}"); return jsonify({"error": str(e)}), 500

@app.route('/api/fetch-image', methods=['POST'])
def api_fetch_image():
    try:
        data = request.json; query = (data.get('query') or '').strip()
        if not query: return jsonify({"error": "Query required"}), 400
        url = fetch_image_for_slide(query, int(data.get('slide_index', 0)), data.get('presentation_id'))
        return jsonify({"success": True, "imageUrl": url})
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/api/update-slide/<pid>/<int:sidx>', methods=['PUT'])
def update_slide(pid, sidx):
    try:
        data = request.json
        if pid not in presentations: return jsonify({"error": "Presentation not found"}), 404
        pres  = presentations[pid]
        if sidx < 0 or sidx >= len(pres["slides"]): return jsonify({"error": "Invalid index"}), 400
        slide = pres["slides"][sidx]
        for k in ("title", "content", "imageUrl"):
            if k in data: slide[k] = data[k]
        return jsonify({"success": True, "slide": slide})
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/api/add-slide/<pid>', methods=['POST'])
def add_slide(pid):
    try:
        if pid not in presentations: return jsonify({"error": "Presentation not found"}), 404
        pres = presentations[pid]; idx = len(pres["slides"])
        url  = fetch_image_for_slide("presentation", idx, pid)
        new  = {"title": "New Slide", "content": ["Point 1", "Point 2", "Point 3"],
                "imageUrl": url, "imageKeywords": "presentation"}
        pres["slides"].append(new)
        return jsonify({"success": True, "slide_index": idx, "slide": new})
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/api/delete-slide/<pid>/<int:sidx>', methods=['DELETE'])
def delete_slide(pid, sidx):
    try:
        if pid not in presentations: return jsonify({"error": "Presentation not found"}), 404
        pres = presentations[pid]
        if sidx < 0 or sidx >= len(pres["slides"]): return jsonify({"error": "Invalid index"}), 400
        if len(pres["slides"]) <= 1: return jsonify({"error": "Cannot delete last slide"}), 400
        pres["slides"].pop(sidx)
        return jsonify({"success": True})
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/api/reorder-slides/<pid>', methods=['PUT'])
def reorder_slides(pid):
    try:
        data = request.json; order = data.get('order', [])
        if pid not in presentations: return jsonify({"error": "Presentation not found"}), 404
        pres = presentations[pid]
        if len(order) != len(pres["slides"]): return jsonify({"error": "Invalid order"}), 400
        pres["slides"] = [pres["slides"][i] for i in order]
        return jsonify({"success": True, "slides": pres["slides"]})
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/api/download-pdf/<pid>', methods=['GET'])
def download_pdf(pid):
    if not AI_ENABLED: return jsonify({"error": "Install reportlab and Pillow for PDF export"}), 503
    try:
        if pid not in presentations: return jsonify({"error": "Presentation not found"}), 404
        pres = presentations[pid]; buf = BytesIO()
        pdf  = pdf_canvas.Canvas(buf, pagesize=landscape(letter))
        W, H = landscape(letter)
        for idx, slide in enumerate(pres["slides"]):
            pdf.setFillColorRGB(0.97,0.97,1.0); pdf.rect(0,0,W,H,fill=True,stroke=False)
            pdf.setFillColorRGB(0.39,0.40,0.95); pdf.rect(0,H-8,W,8,fill=True,stroke=False)
            pdf.setFont("Helvetica-Bold",32); pdf.setFillColorRGB(0.1,0.1,0.2)
            pdf.drawString(50,H-55,slide["title"][:55])
            pdf.setStrokeColorRGB(0.8,0.8,0.95); pdf.line(50,H-68,W-50,H-68)
            pdf.setFont("Helvetica",12); pdf.setFillColorRGB(0.2,0.2,0.35)
            y = H-95
            for pt in slide.get("content",[])[:6]:
                pdf.drawString(50,y,f"  \u2022  {pt[:95]}"); y -= 25
            try:
                rr = requests.get(slide.get("imageUrl",""), timeout=8)
                if rr.status_code == 200:
                    img = Image.open(io.BytesIO(rr.content))
                    img.thumbnail((350,240), Image.Resampling.LANCZOS)
                    ibuf = io.BytesIO(); img.save(ibuf, format="PNG"); ibuf.seek(0)
                    pdf.drawImage(ImageReader(ibuf), W-390, H-310, width=340, height=220)
            except: pass
            pdf.setFont("Helvetica",9); pdf.setFillColorRGB(0.5,0.5,0.6)
            pdf.drawRightString(W-30,18,f"Slide {idx+1} / {len(pres['slides'])}")
            if idx < len(pres["slides"])-1: pdf.showPage()
        pdf.save(); buf.seek(0)
        safe = pres["topic"].replace(" ","_")[:40]
        return send_file(buf, mimetype="application/pdf", as_attachment=True,
                         download_name=f"{safe}_presentation.pdf")
    except Exception as e:
        print(f"PDF error: {e}"); return jsonify({"error": str(e)}), 500

@app.route('/api/download-ppt/<pid>', methods=['GET'])
def download_ppt(pid):
    if not AI_ENABLED:
        return jsonify({"error": "Install python-pptx for PPT export"}), 503
    try:
        if pid not in presentations:
            return jsonify({"error": "Presentation not found"}), 404

        pres  = presentations[pid]
        prs   = PPTPresentation()

        # Slide size — widescreen 16:9
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Theme colors
        BG_COLOR    = RGBColor(0xEE, 0xF7, 0xF6)   # teal light
        TITLE_COLOR = RGBColor(0x1A, 0x33, 0x33)
        BODY_COLOR  = RGBColor(0x2A, 0x60, 0x60)
        ACCENT      = RGBColor(0x2A, 0xA6, 0xA6)

        for slide_data in pres["slides"]:
            slide_layout = prs.slide_layouts[6]  # blank layout
            slide = prs.slides.add_slide(slide_layout)

            # ── Background ──
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = BG_COLOR

            # ── Accent bar at top ──
            bar = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(0), Inches(0),
                Inches(13.33), Inches(0.12)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = ACCENT
            bar.line.fill.background()

            # ── Title ──
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.25),
                Inches(8.5), Inches(0.9)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", "")[:60]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = TITLE_COLOR

            # ── Divider line ──
            line = slide.shapes.add_shape(
                1,
                Inches(0.5), Inches(1.2),
                Inches(8.3), Inches(0.02)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = ACCENT
            line.line.fill.background()

            # ── Bullet points ──
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3),
                Inches(8.0), Inches(5.5)
            )
            tf2 = content_box.text_frame
            tf2.word_wrap = True
            for i, point in enumerate(slide_data.get("content", [])[:6]):
                p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                p2.text = f"• {point[:100]}"
                p2.font.size = Pt(14)
                p2.font.color.rgb = BODY_COLOR
                p2.space_after = Pt(8)

            # ── Image (right side) ──
            img_url = slide_data.get("imageUrl", "")
            if img_url and AI_ENABLED:
                try:
                    r = requests.get(img_url, timeout=8)
                    if r.status_code == 200:
                        img_stream = io.BytesIO(r.content)
                        slide.shapes.add_picture(
                            img_stream,
                            Inches(9.2), Inches(1.2),
                            Inches(3.8), Inches(2.8)
                        )
                except Exception as img_e:
                    print(f"PPT image error: {img_e}")

        # Save to buffer
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        safe = pres["topic"].replace(" ", "_")[:40]
        return send_file(
            buf,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=f"{safe}_presentation.pptx"
        )
    except Exception as e:
        print(f"PPT error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/editorg/<pid>')
def editorg(pid):
    if pid not in presentations: return "Presentation not found", 404
    return render_template('editorg.html', presentation=presentations[pid])

# =============================================================================
#  ERROR HANDLERS & MAIN
# =============================================================================

@app.errorhandler(413)
def too_large(e):
    return jsonify({'error': 'File too large (max 16MB)'}), 413

if __name__ == "__main__":
    init_db()
    print("\n✦ DesignEase Server — http://localhost:8000")
    print("  Admin: admin@designease.com / admin@1156\n")
    app.run(port=8000, debug=True)